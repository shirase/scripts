#!/usr/bin/env python3

"""
convert_to_md.py — Конвертация документов в Markdown
Поддерживаемые форматы: .docx, .doc, .xlsx, .xls, .ods, .pdf, .pptx, .ppt, .csv, .tsv, .html, .epub, .odt, .rtf, .txt

Зависимости:
    pip install pypdf pandas openpyxl python-pptx markdownify odfpy
    sudo apt install pandoc          # или brew install pandoc
    sudo apt install libreoffice     # для .doc / .ppt (legacy)

Использование:
    python convert_to_md.py file.docx
    python convert_to_md.py file.pdf -o result.md
    python convert_to_md.py ./docs/   # конвертировать всю папку
    python convert_to_md.py *.xlsx --output-dir ./markdown/
"""

import sys
import argparse
import subprocess
import shutil
from pathlib import Path


# ── Вспомогательные функции ──────────────────────────────────────────────────

def run(cmd: list[str], **kwargs) -> subprocess.CompletedProcess:
    return subprocess.run(cmd, capture_output=True, text=True, **kwargs)


def pandoc_available() -> bool:
    return shutil.which("pandoc") is not None


def soffice_available() -> bool:
    return shutil.which("soffice") is not None


def save_md(text: str, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(text.strip() + "\n", encoding="utf-8")
    print(f"  ✓  Сохранено: {output_path}")


def html_to_markdown(html: str) -> str:
    """HTML строка в Markdown через markdownify, иначе pandoc."""
    try:
        from markdownify import markdownify as md
        return md(html, heading_style="ATX")
    except ImportError:
        pass

    if pandoc_available():
        result = run(["pandoc", "-f", "html", "-t", "markdown", "--wrap=none"], input=html)
        if result.returncode == 0:
            return result.stdout

    raise RuntimeError(
        "Не удалось сконвертировать HTML в Markdown: установите markdownify или pandoc"
    )


def extract_html_from_mime(src: Path) -> str | None:
    """Если файл содержит MIME multipart с HTML, вернуть HTML содержимое."""
    try:
        from email import policy
        from email.parser import BytesParser
    except ImportError:
        return None

    try:
        data = src.read_bytes()
        msg = BytesParser(policy=policy.default).parsebytes(data)
    except Exception:
        return None

    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/html":
                content = part.get_content()
                if isinstance(content, bytes):
                    content = content.decode(part.get_content_charset("utf-8"), errors="replace")
                return content
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                content = part.get_content()
                if isinstance(content, bytes):
                    content = content.decode(part.get_content_charset("utf-8"), errors="replace")
                return content
    else:
        if msg.get_content_type() == "text/html":
            content = msg.get_content()
            if isinstance(content, bytes):
                content = content.decode(msg.get_content_charset("utf-8"), errors="replace")
            return content

    return None


# ── Конвертеры ───────────────────────────────────────────────────────────────

def convert_docx(src: Path) -> str:
    """DOCX → Markdown через pandoc."""
    if not pandoc_available():
        raise RuntimeError("pandoc не установлен. Установите: sudo apt install pandoc")
    result = run(["pandoc", str(src), "-t", "markdown", "--wrap=none"])
    if result.returncode != 0:
        raise RuntimeError(result.stderr)
    return result.stdout


def convert_doc(src: Path) -> str:
    """DOC (legacy) → DOCX через LibreOffice, затем pandoc.

    В случае MIME-формата (когда .doc фактически контейнер email/multipart с text/html),
    извлекаем HTML и конвертируем напрямую.
    """
    html = extract_html_from_mime(src)
    if html:
        print("  →  Обнаружен MIME-файл с HTML; конвертация html прямо в markdown")
        return html_to_markdown(html)

    if not soffice_available():
        raise RuntimeError(
            "LibreOffice не установлен (нужен для .doc). "
            "Установите: sudo apt install libreoffice"
        )
    import tempfile, os
    with tempfile.TemporaryDirectory() as tmp:
        result = run([
            "soffice", "--headless", "--convert-to", "docx",
            "--outdir", tmp, str(src)
        ])
        if result.returncode != 0:
            raise RuntimeError(result.stderr)
        docx_files = list(Path(tmp).glob("*.docx"))
        if not docx_files:
            raise RuntimeError("LibreOffice не создал .docx файл")
        return convert_docx(docx_files[0])


def convert_xlsx(src: Path) -> str:
    """XLSX / XLSM → Markdown таблицы через openpyxl + pandas."""
    import pandas as pd
    xl = pd.ExcelFile(str(src), engine="openpyxl")
    parts = [f"# {src.stem}\n"]
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        if df.empty:
            continue
        parts.append(f"## {sheet}\n")
        parts.append(df.to_markdown(index=False))
        parts.append("")
    return "\n".join(parts)


def convert_xls(src: Path) -> str:
    """XLS (legacy) → Markdown через pandas + xlrd."""
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(src), engine="xlrd")
    except ImportError:
        raise RuntimeError(
            "xlrd не установлен (нужен для .xls). Установите: pip install xlrd"
        )
    parts = [f"# {src.stem}\n"]
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        if df.empty:
            continue
        parts.append(f"## {sheet}\n")
        parts.append(df.to_markdown(index=False))
        parts.append("")
    return "\n".join(parts)


def convert_ods(src: Path) -> str:
    """ODS → Markdown через pandas + odf."""
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(src), engine="odf")
    except ImportError:
        raise RuntimeError(
            "odfpy не установлен (нужен для .ods). Установите: pip install odfpy"
        )
    parts = [f"# {src.stem}\n"]
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        if df.empty:
            continue
        parts.append(f"## {sheet}\n")
        parts.append(df.to_markdown(index=False))
        parts.append("")
    return "\n".join(parts)


def convert_csv(src: Path, sep: str = ",") -> str:
    """CSV / TSV → Markdown таблица."""
    import pandas as pd
    df = pd.read_csv(str(src), sep=sep)
    return f"# {src.stem}\n\n{df.to_markdown(index=False)}\n"


def convert_pdf(src: Path) -> str:
    """PDF → Markdown (текстовые PDF через pypdf, pandoc как запасной вариант)."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(src))
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"## Страница {i}\n\n{text.strip()}")
        if pages:
            return f"# {src.stem}\n\n" + "\n\n---\n\n".join(pages)
    except Exception:
        pass

    # Запасной вариант — pandoc
    if pandoc_available():
        result = run(["pandoc", str(src), "-t", "markdown", "--wrap=none"])
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout

    raise RuntimeError(
        "Не удалось извлечь текст из PDF. "
        "Для сканированных PDF нужен OCR (tesseract)."
    )


def convert_pptx(src: Path) -> str:
    """PPTX → Markdown (заголовки + текст слайдов)."""
    from pptx import Presentation
    prs = Presentation(str(src))
    parts = [f"# {src.stem}\n"]
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    slide_texts.append(line)
        if slide_texts:
            parts.append(f"## Слайд {i}\n")
            # Первый непустой текст — вероятно заголовок
            parts.append(f"**{slide_texts[0]}**\n" if len(slide_texts) > 1 else slide_texts[0])
            for line in slide_texts[1:]:
                parts.append(f"- {line}")
            parts.append("")
    return "\n".join(parts)


def convert_ppt(src: Path) -> str:
    """PPT (legacy) → PPTX через LibreOffice, затем python-pptx."""
    if not soffice_available():
        raise RuntimeError(
            "LibreOffice не установлен (нужен для .ppt). "
            "Установите: sudo apt install libreoffice"
        )
    import tempfile
    with tempfile.TemporaryDirectory() as tmp:
        result = run([
            "soffice", "--headless", "--convert-to", "pptx",
            "--outdir", tmp, str(src)
        ])
        if result.returncode != 0:
            raise RuntimeError(result.stderr)
        pptx_files = list(Path(tmp).glob("*.pptx"))
        if not pptx_files:
            raise RuntimeError("LibreOffice не создал .pptx файл")
        return convert_pptx(pptx_files[0])


def convert_html(src: Path) -> str:
    """HTML → Markdown через markdownify."""
    try:
        from markdownify import markdownify as md
        html = src.read_text(encoding="utf-8", errors="replace")
        return md(html, heading_style="ATX")
    except ImportError:
        pass
    # Запасной вариант — pandoc
    if pandoc_available():
        result = run(["pandoc", str(src), "-t", "markdown", "--wrap=none"])
        if result.returncode == 0:
            return result.stdout
    raise RuntimeError("markdownify не установлен. pip install markdownify")


def convert_via_pandoc(src: Path) -> str:
    """Универсальная конвертация через pandoc (epub, odt, rtf и др.)."""
    if not pandoc_available():
        raise RuntimeError("pandoc не установлен. Установите: sudo apt install pandoc")
    result = run(["pandoc", str(src), "-t", "markdown", "--wrap=none"])
    if result.returncode != 0:
        raise RuntimeError(result.stderr)
    return result.stdout


def convert_txt(src: Path) -> str:
    """TXT → Markdown (обёртка с заголовком)."""
    text = src.read_text(encoding="utf-8", errors="replace")
    return f"# {src.stem}\n\n{text}"


# ── Диспетчер ─────────────────────────────────────────────────────────────────

CONVERTERS = {
    ".docx": convert_docx,
    ".doc":  convert_doc,
    ".xlsx": convert_xlsx,
    ".xlsm": convert_xlsx,
    ".xls":  convert_xls,
    ".ods":  convert_ods,
    ".csv":  convert_csv,
    ".tsv":  lambda p: convert_csv(p, sep="\t"),
    ".pdf":  convert_pdf,
    ".pptx": convert_pptx,
    ".ppt":  convert_ppt,
    ".html": convert_html,
    ".htm":  convert_html,
    ".epub": convert_via_pandoc,
    ".odt":  convert_via_pandoc,
    ".rtf":  convert_via_pandoc,
    ".txt":  convert_txt,
    ".md":   convert_txt,
}


def convert_file(src: Path, output_dir: Path | None, output_file: Path | None) -> bool:
    ext = src.suffix.lower()
    if ext not in CONVERTERS:
        print(f"  ✗  Неизвестный формат: {src.name}")
        return False

    print(f"  →  Конвертирую: {src.name}  ({ext})")
    try:
        md_text = CONVERTERS[ext](src)
    except Exception as e:
        print(f"  ✗  Ошибка: {e}")
        return False

    if output_file:
        dest = output_file
    elif output_dir:
        dest = output_dir / src.with_suffix(".md").name
    else:
        dest = src.with_suffix(".md")

    save_md(md_text, dest)
    return True


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Конвертация документов в Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "inputs", nargs="+", metavar="FILE_OR_DIR",
        help="Файлы или папки для конвертации",
    )
    parser.add_argument(
        "-o", "--output", metavar="OUTPUT.md",
        help="Выходной файл (только для одного входного файла)",
    )
    parser.add_argument(
        "--output-dir", metavar="DIR",
        help="Папка для сохранения .md файлов",
    )

    args = parser.parse_args()

    output_file = Path(args.output) if args.output else None
    output_dir  = Path(args.output_dir) if args.output_dir else None

    # Собираем список файлов
    files: list[Path] = []
    for inp in args.inputs:
        p = Path(inp)
        if p.is_dir():
            for ext in CONVERTERS:
                files.extend(p.glob(f"*{ext}"))
                files.extend(p.glob(f"*{ext.upper()}"))
        elif "*" in inp or "?" in inp:
            from glob import glob
            files.extend(Path(f) for f in glob(inp))
        elif p.exists():
            files.append(p)
        else:
            print(f"Файл не найден: {inp}", file=sys.stderr)

    if not files:
        print("Не найдено файлов для конвертации.", file=sys.stderr)
        sys.exit(1)

    if output_file and len(files) > 1:
        print("Флаг -o можно использовать только с одним файлом.", file=sys.stderr)
        sys.exit(1)

    ok = sum(convert_file(f, output_dir, output_file if len(files) == 1 else None) for f in files)
    print(f"\nГотово: {ok}/{len(files)} файлов успешно конвертировано.")


if __name__ == "__main__":
    main()
